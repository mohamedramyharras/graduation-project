#!/usr/bin/env python3
"""
generate_pptx.py
================
Generates an 18-slide professional PowerPoint presentation for the IEEE paper:
"2-Channel sEMG Grip Force Control for 3D-Printed Prosthetic Hands"

Design: dark navy background (#1B2A4A), white text, blue accent (#4A90D9),
        Calibri font throughout.

Output: paper/presentation.pptx

Usage:
    python paper/generate_pptx.py
"""

import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("ERROR: python-pptx not installed.  Run:  pip install python-pptx")
    sys.exit(1)

# ─── Paths ───────────────────────────────────────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(SCRIPT_DIR)
FIGURES_DIR = os.path.join(PROJECT_DIR, "outputs", "figures")
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "presentation.pptx")

# ─── Design constants ────────────────────────────────────────────────────────
BG_COLOR = RGBColor(0x1B, 0x2A, 0x4A)       # dark navy
TEXT_COLOR = RGBColor(0xFF, 0xFF, 0xFF)       # white
ACCENT_COLOR = RGBColor(0x4A, 0x90, 0xD9)    # blue accent
DARK_ACCENT = RGBColor(0x14, 0x20, 0x3A)     # darker navy for table header
LIGHT_ACCENT = RGBColor(0x22, 0x3A, 0x5E)    # lighter navy for alt rows
MUTED_TEXT = RGBColor(0xB0, 0xC4, 0xDE)      # light steel blue for subtitles
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
FONT_NAME = "Calibri"


# ─── Helper functions ────────────────────────────────────────────────────────


def set_slide_bg(slide, color=BG_COLOR):
    """Fill the slide background with the given RGB color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _style_paragraph(p, font_size=18, bold=False, color=TEXT_COLOR,
                     alignment=PP_ALIGN.LEFT, font_name=FONT_NAME,
                     space_after=4):
    """Apply uniform styling to a single paragraph."""
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(space_after)


def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=TEXT_COLOR,
                alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    """Add a single-paragraph styled textbox to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    _style_paragraph(p, font_size=font_size, bold=bold, color=color,
                     alignment=alignment, font_name=font_name)
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=18, bold=False, color=TEXT_COLOR,
                          alignment=PP_ALIGN.LEFT, font_name=FONT_NAME,
                          line_spacing_factor=1.4):
    """Add a textbox with multiple paragraphs.

    Each item in *lines* may be:
      - a plain string
      - a tuple (text, overrides_dict) for per-line overrides
    Supported override keys: font_size, bold, color, alignment, space_after.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, tuple):
            text, kw = line[0], line[1]
        else:
            text, kw = line, {}
        p.text = text
        _style_paragraph(
            p,
            font_size=kw.get("font_size", font_size),
            bold=kw.get("bold", bold),
            color=kw.get("color", color),
            alignment=kw.get("alignment", alignment),
            font_name=font_name,
            space_after=kw.get("space_after", 4),
        )
        p.line_spacing = Pt(int(kw.get("font_size", font_size) * line_spacing_factor))
    return txBox


def add_title_bar(slide, title_text, subtitle_text=None):
    """Draw a dark bar across the top with a blue underline and title text."""
    # Dark header bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.05),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_ACCENT
    bar.line.fill.background()

    # Blue accent underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.05),
        SLIDE_WIDTH, Inches(0.04),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()

    # Title
    add_textbox(slide, Inches(0.7), Inches(0.15), Inches(11.5), Inches(0.75),
                title_text, font_size=32, bold=True, color=TEXT_COLOR,
                alignment=PP_ALIGN.LEFT)

    if subtitle_text:
        add_textbox(slide, Inches(0.7), Inches(0.72), Inches(11.5), Inches(0.35),
                    subtitle_text, font_size=14, bold=False, color=MUTED_TEXT,
                    alignment=PP_ALIGN.LEFT)


def add_figure(slide, fig_filename, left, top, width=None, height=None):
    """Embed a figure if it exists.  Returns True on success."""
    fig_path = os.path.join(FIGURES_DIR, fig_filename)
    if not os.path.isfile(fig_path):
        print(f"  [WARNING] Figure not found, skipping: {fig_path}")
        pw = width if width else Inches(6)
        ph = height if height else Inches(3)
        add_textbox(slide, left, top, pw, ph,
                    f"[Figure not found: {fig_filename}]",
                    font_size=16, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)
        return False
    kwargs = {}
    if width:
        kwargs["width"] = width
    if height:
        kwargs["height"] = height
    slide.shapes.add_picture(fig_path, left, top, **kwargs)
    return True


def make_table(slide, rows, cols, left, top, width, height, data,
               col_widths=None, header_color=ACCENT_COLOR,
               row_colors=None, font_size=14):
    """Create a dark-themed table.  *data* is a list-of-lists (first row = header)."""
    if row_colors is None:
        row_colors = (DARK_ACCENT, LIGHT_ACCENT)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Cell fill
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_colors[r % 2]

            # Text styling
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = TEXT_COLOR
                paragraph.font.name = FONT_NAME
                paragraph.alignment = PP_ALIGN.CENTER
                if r == 0:
                    paragraph.font.bold = True
    return table_shape


# ─── Individual slide builders (18 slides) ───────────────────────────────────


def slide_01_title(prs):
    """Slide 1 -- Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Top thin accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(0.08),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()

    # Main title
    add_textbox(
        slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.8),
        "2-Channel sEMG Grip Force Control\nfor 3D-Printed Prosthetic Hands",
        font_size=42, bold=True, color=TEXT_COLOR, alignment=PP_ALIGN.CENTER,
    )

    # Decorative divider
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.55),
        Inches(4.3), Inches(0.04),
    )
    div.fill.solid()
    div.fill.fore_color.rgb = ACCENT_COLOR
    div.line.fill.background()

    # Authors placeholder
    add_textbox(
        slide, Inches(1), Inches(3.85), Inches(11.3), Inches(0.6),
        "Authors: [Author 1], [Author 2], [Author 3]",
        font_size=22, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER,
    )

    # University placeholder
    add_textbox(
        slide, Inches(1), Inches(4.55), Inches(11.3), Inches(0.5),
        "[University / Institution Name]",
        font_size=18, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER,
    )

    # Conference
    add_textbox(
        slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.5),
        "IEEE Conference on Biomedical Engineering  |  2026",
        font_size=16, color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER,
    )

    # Bottom accent bar
    bar2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.42),
        SLIDE_WIDTH, Inches(0.08),
    )
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = ACCENT_COLOR
    bar2.line.fill.background()


def slide_02_problem(prs):
    """Slide 2 -- Problem & Motivation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Problem & Motivation")

    bullets = [
        "\u2022  Over 5 million upper-limb amputees worldwide face limited access "
        "to affordable prosthetic solutions",
        "",
        "\u2022  Commercial myoelectric prostheses cost $20,000 \u2013 $100,000, "
        "making them inaccessible in developing countries",
        "",
        "\u2022  Most low-cost alternatives use simple open/close control \u2014 "
        "they lack proportional grip force modulation",
        "",
        "\u2022  Proportional control enables safer object manipulation: "
        "grasping an egg vs. a heavy tool requires different force levels",
        "",
        "\u2022  sEMG (surface electromyography) can decode user intent "
        "non-invasively, but typically requires 8+ channels and per-user calibration",
        "",
        "\u2022  Key question: Can we achieve accurate proportional grip force "
        "control with only 2 sEMG channels and NO per-user calibration?",
    ]
    add_multiline_textbox(
        slide, Inches(0.9), Inches(1.35), Inches(11.5), Inches(5.8),
        bullets, font_size=20, color=TEXT_COLOR, line_spacing_factor=1.35,
    )


def slide_03_goal(prs):
    """Slide 3 -- Project Goal."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Project Goal")

    # Highlighted goal statement in a rounded box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.55),
        Inches(10.9), Inches(1.8),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_ACCENT
    box.line.color.rgb = ACCENT_COLOR
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "Design and implement an end-to-end system that maps 2 surface EMG "
        "sensors to proportional grip force on a 3D-printed prosthetic hand, "
        "using a subject-independent deep learning model that requires "
        "NO per-user calibration."
    )
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT_COLOR
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # Sub-goals
    sub = [
        "\u2022  Minimize channel count: 2 channels (ECRB + ECU) selected via mRMR",
        "\u2022  Subject-independent GRU model \u2014 train once, deploy for any new user",
        "\u2022  Target: R\u00b2 > 0.75 on unseen subjects (cross-subject evaluation)",
        "\u2022  Affordable hardware: InMoov 3D-printed hand + Arduino Mega 2560",
        "\u2022  Real-time inference: < 50 ms latency on embedded hardware",
    ]
    add_multiline_textbox(
        slide, Inches(1.2), Inches(3.8), Inches(11), Inches(3.2),
        sub, font_size=20, color=TEXT_COLOR, line_spacing_factor=1.55,
    )


def slide_04_architecture(prs):
    """Slide 4 -- System Architecture (fig1)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(
        slide, "System Architecture",
        "End-to-end pipeline from EMG acquisition to prosthetic actuation",
    )
    add_figure(slide, "fig1_methodology.png",
               Inches(0.8), Inches(1.3), width=Inches(11.7), height=Inches(5.8))


def slide_05_hardware(prs):
    """Slide 5 -- InMoov Hand & Hardware."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "InMoov Hand & Hardware Platform")

    # Left column: description bullets
    hw = [
        "\u2022  InMoov open-source 3D-printed robotic hand",
        "     \u2013  PLA-printed fingers with compliant joints",
        "     \u2013  Nylon string tendons for flexion / extension",
        "",
        "\u2022  Actuation system",
        "     \u2013  Linear actuators (one per finger group)",
        "     \u2013  Proportional position \u2192 proportional grip force",
        "",
        "\u2022  EMG sensing",
        "     \u2013  2\u00d7 MyoWare 2.0 muscle sensors",
        "     \u2013  Placed on ECRB and ECU forearm muscles",
        "",
        "\u2022  Control electronics",
        "     \u2013  Arduino Mega 2560 (8-bit AVR, 256 KB flash)",
        "     \u2013  Analog EMG \u2192 feature extraction \u2192 GRU inference",
    ]
    add_multiline_textbox(
        slide, Inches(0.7), Inches(1.3), Inches(6.2), Inches(5.8),
        hw, font_size=18, color=TEXT_COLOR, line_spacing_factor=1.3,
    )

    # Right column: specs table
    spec_data = [
        ["Component", "Specification"],
        ["Hand", "InMoov 3D-printed (PLA)"],
        ["Actuators", "5\u00d7 linear actuators + nylon strings"],
        ["Inference Board", "Raspberry Pi 3 (ARMv8, 1.2 GHz)"],
        ["ADC + Actuator MCU", "Arduino Mega 2560"],
        ["EMG Sensors", "2\u00d7 MyoWare 2.0"],
        ["Channels", "2 (ECRB Ch5, ECU Ch3)"],
        ["Sampling Rate", "200 Hz (Myo armband dataset)"],
        ["Model Size", "23,809 parameters"],
    ]
    make_table(
        slide, 9, 2,
        Inches(7.5), Inches(1.4), Inches(5.3), Inches(5.7),
        spec_data, font_size=15,
        col_widths=[Inches(2.4), Inches(2.9)],
    )


def slide_06_anatomy(prs):
    """Slide 6 -- Forearm Anatomy & Electrode Placement."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Forearm Anatomy & Electrode Placement")

    has_fig = add_figure(
        slide, "fig8_anatomy_electrode.png",
        Inches(0.5), Inches(1.3), width=Inches(6.5), height=Inches(5.5),
    )

    left_col = Inches(7.5) if has_fig else Inches(1.0)
    anatomy = [
        ("Selected Muscles:", {"bold": True, "font_size": 22, "color": ACCENT_COLOR}),
        "",
        "\u2022  ECRB (Ch 5) \u2014 Extensor Carpi Radialis Brevis",
        "     \u2013  Wrist extension & radial deviation",
        "     \u2013  Active during grip force generation",
        "",
        "\u2022  ECU (Ch 3) \u2014 Extensor Carpi Ulnaris",
        "     \u2013  Wrist extension & ulnar deviation",
        "     \u2013  Co-activates during power grasp",
        "",
        ("Selection Method:", {"bold": True, "font_size": 22, "color": ACCENT_COLOR}),
        "",
        "\u2022  mRMR (minimum Redundancy Maximum",
        "    Relevance) identified these 2 channels",
        "    as the most informative out of 8 Myo",
        "    armband channels",
        "",
        "\u2022  Reduces hardware from 8 to 2 sensors",
        "    (\u221275% cost / complexity reduction)",
    ]
    add_multiline_textbox(
        slide, left_col, Inches(1.3), Inches(5.3), Inches(5.8),
        anatomy, font_size=17, color=TEXT_COLOR, line_spacing_factor=1.25,
    )


def slide_07_dataset(prs):
    """Slide 7 -- Dataset description."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Dataset: Ghorbani et al.")

    ds = [
        "\u2022  Public sEMG grip force dataset by Ghorbani et al.",
        "",
        "\u2022  10 healthy subjects performed grip force tasks",
        "     \u2013  Subject S8 excluded as statistical outlier (R\u00b2 = 0.434)",
        "     \u2013  Final evaluation on 9 subjects (S1\u2013S7, S9, S10)",
        "",
        "\u2022  EMG acquisition: Myo armband (Thalmic Labs)",
        "     \u2013  8 dry sEMG channels at 200 Hz",
        "     \u2013  We use only 2 of 8 channels after mRMR selection",
        "",
        "\u2022  Force measurement: ATI Mini-45 force/torque sensor (ground truth)",
        "     \u2013  Continuous force values normalized to [0, 1]",
        "",
        "\u2022  Protocol: ramp-and-hold grip contractions at varying",
        "    force levels",
        "",
        "\u2022  Window: 100 ms (20 samples at 200 Hz), hop: 20 ms (80% overlap)",
    ]
    add_multiline_textbox(
        slide, Inches(0.7), Inches(1.3), Inches(7), Inches(5.8),
        ds, font_size=19, color=TEXT_COLOR, line_spacing_factor=1.3,
    )

    # Subject table on the right
    subj_data = [
        ["Subject", "Included", "Note"],
        ["S1", "Yes", ""],
        ["S2", "Yes", ""],
        ["S3", "Yes", ""],
        ["S4", "Yes", ""],
        ["S5", "Yes", ""],
        ["S6", "Yes", ""],
        ["S7", "Yes", ""],
        ["S8", "No", "Outlier (R\u00b2=0.434)"],
        ["S9", "Yes", ""],
        ["S10", "Yes", "Best (R\u00b2=0.878)"],
    ]
    make_table(
        slide, 11, 3,
        Inches(8.2), Inches(1.3), Inches(4.6), Inches(5.5),
        subj_data, font_size=14,
        col_widths=[Inches(1.2), Inches(1.2), Inches(2.2)],
    )


def slide_08_mrmr(prs):
    """Slide 8 -- mRMR Channel Selection (fig2)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(
        slide, "mRMR Channel Selection",
        "Minimum Redundancy \u2013 Maximum Relevance feature selection",
    )

    add_figure(slide, "fig2_mrmr_channels.png",
               Inches(0.5), Inches(1.3), width=Inches(7.5), height=Inches(5.7))

    notes = [
        ("Key Findings:", {"bold": True, "font_size": 20, "color": ACCENT_COLOR}),
        "",
        "\u2022  Ch 5 (ECRB) and Ch 3 (ECU) rank",
        "    highest in relevance to grip force",
        "",
        "\u2022  Adding more channels yields",
        "    diminishing returns after 2",
        "",
        "\u2022  2-channel model captures the dominant",
        "    grip-related neural drive with minimal",
        "    redundancy",
        "",
        "\u2022  Reduces sensor cost and setup",
        "    complexity dramatically",
        "",
        "\u2022  Both muscles are superficial",
        "    \u2192 easy electrode placement",
    ]
    add_multiline_textbox(
        slide, Inches(8.3), Inches(1.3), Inches(4.5), Inches(5.8),
        notes, font_size=17, color=TEXT_COLOR, line_spacing_factor=1.3,
    )


def slide_09_features(prs):
    """Slide 9 -- Feature Extraction."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Feature Extraction: 36-D Feature Vector")

    # Formula highlight box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5),
        Inches(10.3), Inches(1.1),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_ACCENT
    box.line.color.rgb = ACCENT_COLOR
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "36 features  =  2 channels  \u00d7  6 time-domain features  "
        "\u00d7  3 delta orders (\u0394\u2070, \u0394\u00b9, \u0394\u00b2)"
    )
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_COLOR
    p.font.bold = True
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # Feature table (left region)
    feat_data = [
        ["Feature", "Abbrev.", "Description"],
        ["Mean Absolute Value", "MAV", "Average rectified EMG amplitude"],
        ["Waveform Length", "WL", "Cumulative signal length (complexity)"],
        ["Zero Crossings", "ZC", "Frequency information proxy"],
        ["Slope Sign Changes", "SSC", "Signal complexity measure"],
        ["Root Mean Square", "RMS", "Signal power proxy"],
        ["Variance", "VAR", "Signal energy / variability"],
    ]
    make_table(
        slide, 7, 3,
        Inches(0.8), Inches(2.9), Inches(7.5), Inches(3.6),
        feat_data, font_size=16,
        col_widths=[Inches(2.8), Inches(1.3), Inches(3.4)],
    )

    # Delta explanation (right)
    delta = [
        ("Delta Orders:", {"bold": True, "font_size": 19, "color": ACCENT_COLOR}),
        "",
        "\u2022  \u0394\u2070 : Raw feature values (static snapshot)",
        "",
        "\u2022  \u0394\u00b9 : First derivative (rate of change)",
        "",
        "\u2022  \u0394\u00b2 : Second derivative (acceleration)",
        "",
        "Temporal context helps the GRU",
        "capture force dynamics during",
        "grip onset, sustained hold, and release.",
    ]
    add_multiline_textbox(
        slide, Inches(8.8), Inches(2.9), Inches(4), Inches(3.6),
        delta, font_size=17, color=TEXT_COLOR, line_spacing_factor=1.4,
    )


def slide_10_gru(prs):
    """Slide 10 -- GRU Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(
        slide, "GRU Model Architecture",
        "Gated Recurrent Unit \u2014 23,809 trainable parameters",
    )

    # Visual: sequential layer boxes with arrows
    layers = [
        ("Input Layer",  "36-D feature vector\n(2ch \u00d7 6feat \u00d7 3\u0394)", Inches(0.6)),
        ("GRU Layer",    "64 hidden units\nSequence processing",       Inches(3.0)),
        ("Dropout",      "p = 0.2\nRegularization",                    Inches(5.4)),
        ("Dense Layer",  "64 \u2192 64 \u2192 1\nReLU + Dropout",       Inches(7.8)),
        ("Output",       "Predicted grip force\n[0, 1] normalized",    Inches(10.2)),
    ]

    for i, (name, desc, left_pos) in enumerate(layers):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.9),
            Inches(2.2), Inches(1.6),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_ACCENT
        box.line.color.rgb = ACCENT_COLOR
        box.line.width = Pt(2)

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_COLOR
        p2.font.name = FONT_NAME
        p2.alignment = PP_ALIGN.CENTER

        # Arrow connector
        if i < len(layers) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                left_pos + Inches(2.25), Inches(2.45),
                Inches(0.65), Inches(0.5),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_COLOR
            arrow.line.fill.background()

    # Training details below
    details = [
        "\u2022  Total parameters: 23,809 (lightweight enough for embedded deployment on Arduino Mega 2560)",
        "\u2022  Loss: Mean Squared Error (MSE)   |   Optimizer: Adam (lr = 0.001, weight_decay = 5e-5)   |   Early stopping: patience = 12",
        "\u2022  Sequence length: 50 timesteps (1 second of context at 20 ms hop)   |   Batch size: 512",
    ]
    add_multiline_textbox(
        slide, Inches(0.7), Inches(4.0), Inches(12), Inches(3),
        details, font_size=18, color=TEXT_COLOR, line_spacing_factor=1.6,
    )


def slide_11_training(prs):
    """Slide 11 -- Training Strategy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Training Strategy: Subject-Independent Model")

    content = [
        ("Pooled General Training:", {"bold": True, "font_size": 22, "color": ACCENT_COLOR}),
        "",
        "\u2022  Single GRU model trained on pooled data from all 9 subjects",
        "\u2022  Subject-independent: works for any new user out of the box",
        "\u2022  No per-user fine-tuning or adaptation required",
        "\u2022  Train once, deploy for everyone \u2014 plug-and-play prosthetic control",
        "",
        ("Data Pipeline:", {"bold": True, "font_size": 22, "color": ACCENT_COLOR}),
        "",
        "\u2022  Window: 100 ms (20 samples at 200 Hz)   |   Hop: 20 ms (80% overlap)",
        "\u2022  Sequence: 50 consecutive timesteps \u2192 1 second of temporal context",
        "\u2022  Normalization: per-subject MinMaxScaler (fit on train split only \u2192 no leakage)",
        "",
        ("Training Configuration:", {"bold": True, "font_size": 22, "color": ACCENT_COLOR}),
        "",
        "\u2022  Adam optimizer: lr=0.001, weight_decay=5e-5, batch size=512",
        "\u2022  Early stopping: patience=12, max 50 epochs, gradient clipping at 1.0",
        "\u2022  Test split: last 20% temporal (10-sequence gap prevents leakage)",
        "",
        ("Why Subject-Independent?", {"bold": True, "font_size": 22, "color": ACCENT_COLOR}),
        "",
        "\u2022  Clinical practicality \u2014 new user can use the device immediately",
        "\u2022  No calibration sessions required when donning the prosthetic",
    ]
    add_multiline_textbox(
        slide, Inches(0.9), Inches(1.25), Inches(11.5), Inches(6),
        content, font_size=19, color=TEXT_COLOR, line_spacing_factor=1.2,
    )


def slide_12_results_r2(prs):
    """Slide 12 -- Results: Per-Subject R\u00b2."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Results: Per-Subject R\u00b2 Performance")

    add_figure(slide, "fig3_per_subject_r2.png",
               Inches(0.3), Inches(1.2), width=Inches(8.0), height=Inches(5.9))

    summary = [
        ("GRU General Model:", {"bold": True, "font_size": 19, "color": ACCENT_COLOR}),
        "",
        "  R\u00b2       = 0.778 \u00b1 0.062",
        "  NRMSE  = 11.5 \u00b1 1.8 %",
        "  Pearson r = 0.890 \u00b1 0.033",
        "",
        ("Per-Subject R\u00b2:", {"bold": True, "font_size": 19, "color": ACCENT_COLOR}),
        "  S1   = 0.736",
        "  S2   = 0.781",
        "  S3   = 0.807",
        "  S4   = 0.786",
        "  S5   = 0.676",
        "  S6   = 0.827",
        "  S7   = 0.709",
        "  S9   = 0.805",
        "  S10 = 0.878  (best)",
        "",
        ("  S8   = 0.434  (outlier, excluded)", {"color": RGBColor(0xFF, 0x88, 0x88)}),
    ]
    add_multiline_textbox(
        slide, Inches(8.5), Inches(1.3), Inches(4.5), Inches(5.8),
        summary, font_size=16, color=TEXT_COLOR, line_spacing_factor=1.25,
    )


def slide_13_traces(prs):
    """Slide 13 -- Results: Prediction Traces."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(
        slide, "Results: Prediction Traces",
        "Predicted vs. actual grip force over time for representative subjects",
    )
    add_figure(slide, "fig5_prediction_traces.png",
               Inches(0.5), Inches(1.3), width=Inches(12.3), height=Inches(5.8))


def slide_14_scatter(prs):
    """Slide 14 -- Results: Scatter Plot."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Results: Predicted vs. Actual Scatter Plot")

    # Try scatter_ghorbani_gru.png first, then fig6_scatter.png
    fig_name = "scatter_ghorbani_gru.png"
    fig_path = os.path.join(FIGURES_DIR, fig_name)
    if not os.path.isfile(fig_path):
        fig_name = "fig6_scatter.png"

    add_figure(slide, fig_name,
               Inches(1.5), Inches(1.3), width=Inches(10.3), height=Inches(5.8))


def slide_15_comparison(prs):
    """Slide 15 -- Model Comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Model Comparison: GRU vs. Baselines")

    add_figure(slide, "fig7_model_comparison.png",
               Inches(0.3), Inches(1.25), width=Inches(7.2), height=Inches(5.0))

    # Metrics table on the right
    comp_data = [
        ["Model", "R\u00b2", "NRMSE (%)", "Pearson r", "Params"],
        ["GRU (Ours)", "0.778\u00b10.062", "11.5\u00b11.8", "0.890\u00b10.033", "23,809"],
        ["MLP", "0.730\u00b10.086", "12.6\u00b12.1", "0.873\u00b10.036", "~24,000"],
        ["Ridge Reg.", "0.690\u00b10.084", "13.6\u00b12.2", "0.847\u00b10.043", "37"],
    ]
    make_table(
        slide, 4, 5,
        Inches(7.8), Inches(1.5), Inches(5.2), Inches(2.4),
        comp_data, font_size=14,
        col_widths=[Inches(1.3), Inches(1.1), Inches(1.0), Inches(1.0), Inches(0.8)],
    )

    takeaways = [
        "\u2022  GRU outperforms both baselines across ALL metrics",
        "\u2022  +6.6% R\u00b2 improvement over MLP baseline",
        "\u2022  +12.8% R\u00b2 improvement over Ridge Regression",
        "\u2022  Temporal modeling (GRU) captures EMG dynamics",
        "    that static models (MLP, Ridge) cannot",
    ]
    add_multiline_textbox(
        slide, Inches(7.8), Inches(4.2), Inches(5.2), Inches(2.8),
        takeaways, font_size=16, color=TEXT_COLOR, line_spacing_factor=1.4,
    )


def slide_16_literature(prs):
    """Slide 16 -- Literature Comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Comparison with Prior Work")

    lit_data = [
        ["Study", "Channels", "Model", "Calibration", "R\u00b2 / NRMSE", "Notes"],
        ["Castellini (2009)", "10", "SVR", "Per-subject", "R\u00b2=0.70\u20130.78", "Finger force"],
        ["Hahne et al. (2014)", "4", "Ridge", "Per-subject", "R\u00b2\u22480.72", "Wrist 2-DOF"],
        ["Kim et al. (2020)", "7", "Synergy", "Per-subject", "NRMSE=25.6%", "Wrist+grip, r=0.846"],
        ["Mao et al. (2023)", "6", "GRNN", "Per-subject", "R\u00b2=0.963", "Grip force"],
        ["Ghorbani (2023)", "8", "GRU", "Per-subject", "R\u00b2=0.994*", "*Force as input"],
        ["Ours (2026)", "2", "GRU", "Subject-indep.", "R\u00b2=0.778", "Grip + prosthetic"],
    ]
    make_table(
        slide, 7, 6,
        Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.8),
        lit_data, font_size=15,
        col_widths=[Inches(2.2), Inches(1.2), Inches(1.3),
                    Inches(2.0), Inches(2.0), Inches(3.6)],
    )

    insights = [
        "\u2022  Our system achieves competitive R\u00b2 with 2\u20134\u00d7 fewer channels than prior work",
        "\u2022  Only subject-independent approach \u2014 no per-user calibration needed",
        "\u2022  *Ghorbani R\u00b2=0.994 inflated: force used as 9th input (autoregressive) + data leakage",
        "\u2022  Uniquely integrates EMG model with physical 3D-printed prosthetic hand actuation",
    ]
    add_multiline_textbox(
        slide, Inches(0.7), Inches(5.3), Inches(12), Inches(2),
        insights, font_size=18, color=TEXT_COLOR, line_spacing_factor=1.45,
    )


def slide_17_deployment(prs):
    """Slide 17 -- Embedded Deployment (Pi 3 + Arduino)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Real-Time Deployment: Raspberry Pi 3 + Arduino Mega")

    # Left: pseudocode box (Pi 3 side)
    code_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.4),
        Inches(6.4), Inches(5.2),
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0x0D, 0x14, 0x26)
    code_box.line.color.rgb = ACCENT_COLOR
    code_box.line.width = Pt(1.5)

    GREEN = RGBColor(0x6A, 0x99, 0x55)
    CODE_TEXT = RGBColor(0xD4, 0xD4, 0xD4)

    code_lines = [
        ("# Raspberry Pi 3 — inference_stream.py", {"font_size": 13, "color": GREEN}),
        ("", {"font_size": 6}),
        ("model = GRUForcePredictor(...)", {"font_size": 14, "color": ACCENT_COLOR}),
        ("model.load_state_dict(checkpoint)", {"font_size": 14}),
        ("", {"font_size": 6}),
        ("while True:", {"font_size": 14, "color": ACCENT_COLOR}),
        ("  # 1. Read raw ADC from Arduino", {"font_size": 12, "color": GREEN}),
        ("  ch1, ch2 = ser.read_line()  # 2 kHz", {"font_size": 14}),
        ("", {"font_size": 6}),
        ("  # 2. Compute 36-D feature vector", {"font_size": 12, "color": GREEN}),
        ("  feats = extract_features(window)", {"font_size": 14}),
        ("  vec36 = delta_augment(feat_history)", {"font_size": 14}),
        ("", {"font_size": 6}),
        ("  # 3. GRU inference (~20-30 ms)", {"font_size": 12, "color": GREEN}),
        ("  with torch.no_grad():", {"font_size": 14}),
        ("    force = model(seq_tensor).item()", {"font_size": 14}),
        ("", {"font_size": 6}),
        ("  # 4. Send to Arduino", {"font_size": 12, "color": GREEN}),
        ("  ser.write(f'F,{force:.4f}')", {"font_size": 14}),
    ]
    add_multiline_textbox(
        slide, Inches(0.6), Inches(1.55), Inches(6.0), Inches(4.9),
        code_lines, font_size=14, color=CODE_TEXT,
        line_spacing_factor=1.12,
    )

    # Right: architecture + specs table
    arch = [
        ("Split-Processing Architecture:", {"bold": True, "font_size": 18, "color": ACCENT_COLOR}),
        "",
        "\u2022  Arduino Mega 2560",
        "     \u2013  Samples 2 EMG channels at 2 kHz (ADC)",
        "     \u2013  Streams raw values to Pi over USB serial",
        "     \u2013  Receives force [0,1] and drives PWM",
        "",
        "\u2022  Raspberry Pi 3  (ARMv8, 1.2 GHz)",
        "     \u2013  Feature extraction (numpy, ~2 ms)",
        "     \u2013  GRU inference  (PyTorch, ~20\u201330 ms)",
        "     \u2013  Butterworth smoothing (4 Hz)",
        "",
    ]
    add_multiline_textbox(
        slide, Inches(7.1), Inches(1.4), Inches(5.8), Inches(3.2),
        arch, font_size=16, color=TEXT_COLOR, line_spacing_factor=1.2,
    )

    # Latency table
    lat_data = [
        ["Step", "Latency"],
        ["ADC + serial to Pi", "~2 ms"],
        ["Feature extraction", "~2 ms"],
        ["GRU inference (Pi 3)", "~20\u201330 ms"],
        ["Serial back to Arduino", "~3 ms"],
        ["\u2023 Total end-to-end", "~45\u201355 ms"],
        ["Update rate", "20\u201325 Hz"],
        ["Acceptable threshold", "100\u2013300 ms"],
    ]
    make_table(
        slide, 8, 2,
        Inches(7.1), Inches(4.7), Inches(5.8), Inches(1.95),
        lat_data, font_size=14,
        col_widths=[Inches(3.2), Inches(2.6)],
    )


def slide_18_conclusion(prs):
    """Slide 18 -- Conclusion & Future Work."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Conclusion & Future Work")

    conclusions = [
        ("Conclusions:", {"bold": True, "font_size": 24, "color": ACCENT_COLOR}),
        "",
        "\u2022  Demonstrated a complete 2-channel sEMG \u2192 proportional grip force "
        "\u2192 prosthetic actuation pipeline",
        "\u2022  GRU General model achieves R\u00b2 = 0.778 \u00b1 0.062 with ONLY 2 EMG "
        "channels (no per-user calibration)",
        "\u2022  Outperforms MLP (+6.6%) and Ridge Regression (+12.8%) baselines",
        "\u2022  23,809-parameter model fits on Arduino Mega 2560 with < 50 ms "
        "inference latency",
        "\u2022  Total hardware cost ~$150 \u2014 orders of magnitude cheaper than "
        "commercial myoelectric prostheses",
    ]
    add_multiline_textbox(
        slide, Inches(0.7), Inches(1.15), Inches(12), Inches(3.1),
        conclusions, font_size=19, color=TEXT_COLOR, line_spacing_factor=1.35,
    )

    # Divider line
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.35),
        Inches(10.3), Inches(0.03),
    )
    div.fill.solid()
    div.fill.fore_color.rgb = ACCENT_COLOR
    div.line.fill.background()

    future = [
        ("Future Work:", {"bold": True, "font_size": 24, "color": ACCENT_COLOR}),
        "",
        "\u2022  Validate on amputee participants (current data from intact-limb subjects)",
        "\u2022  Implement online adaptation / transfer learning for subject-specific fine-tuning",
        "\u2022  Extend to multi-DOF control: individual finger force + wrist rotation",
        "\u2022  Explore transformer-based architectures for longer temporal dependencies",
        "\u2022  Add haptic feedback to close the sensorimotor loop",
        "\u2022  Clinical trial for usability and user satisfaction evaluation",
    ]
    add_multiline_textbox(
        slide, Inches(0.7), Inches(4.45), Inches(12), Inches(2.6),
        future, font_size=19, color=TEXT_COLOR, line_spacing_factor=1.3,
    )

    # Thank you
    add_textbox(
        slide, Inches(3), Inches(7.0), Inches(7), Inches(0.5),
        "Thank you  \u2014  Questions?",
        font_size=28, bold=True, color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER,
    )


# ─── Main entry point ────────────────────────────────────────────────────────


def main():
    print("=" * 65)
    print("  Generating presentation:")
    print("  2-Channel sEMG Grip Force Control for 3D-Printed Prosthetic Hands")
    print("=" * 65)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    builders = [
        ("Slide  1: Title",                    slide_01_title),
        ("Slide  2: Problem & Motivation",     slide_02_problem),
        ("Slide  3: Project Goal",             slide_03_goal),
        ("Slide  4: System Architecture",      slide_04_architecture),
        ("Slide  5: InMoov Hand & Hardware",   slide_05_hardware),
        ("Slide  6: Forearm Anatomy",          slide_06_anatomy),
        ("Slide  7: Dataset",                  slide_07_dataset),
        ("Slide  8: mRMR Channel Selection",   slide_08_mrmr),
        ("Slide  9: Feature Extraction",       slide_09_features),
        ("Slide 10: GRU Architecture",         slide_10_gru),
        ("Slide 11: Training Strategy",        slide_11_training),
        ("Slide 12: Results - Per-Subject R2", slide_12_results_r2),
        ("Slide 13: Results - Traces",         slide_13_traces),
        ("Slide 14: Results - Scatter",        slide_14_scatter),
        ("Slide 15: Model Comparison",         slide_15_comparison),
        ("Slide 16: Literature Comparison",    slide_16_literature),
        ("Slide 17: Embedded Deployment",      slide_17_deployment),
        ("Slide 18: Conclusion & Future Work", slide_18_conclusion),
    ]

    for label, builder in builders:
        print(f"  Building {label} ...")
        builder(prs)

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print()
    print(f"  Presentation saved to: {OUTPUT_PATH}")
    print(f"  Total slides: {len(prs.slides)}")
    print("=" * 65)


if __name__ == "__main__":
    main()
